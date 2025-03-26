from pptx import Presentation
from pptx.util import Inches

# Load the existing template
template_path = "path/to/your/template.pptx"  # Replace with your template file path
prs = Presentation(template_path)

# Function to add a slide with an image using a specific layout
def add_slide_with_image(prs, image_path, title_text="Slide Title", layout_index=1):
    # Use a layout from the template (e.g., layout 1 is often "Title and Content")
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title (assuming the layout has a title placeholder)
    title = slide.shapes.title
    title.text = title_text
    
    # Add the image (adjust positioning based on your template)
    left = Inches(1)    # Position from the left
    top = Inches(1.5)   # Position from the top
    height = Inches(4)  # Height of the image
    slide.shapes.add_picture(image_path, left, top, height=height)

# List of images and their titles (customize this)
images = [
    {"path": "path/to/image1.jpg", "title": "Slide 1: First Image"},
    {"path": "path/to/image2.png", "title": "Slide 2: Second Image"},
    {"path": "path/to/image3.jpeg", "title": "Slide 3: Third Image"},
]

# Add slides with images
for image in images:
    add_slide_with_image(prs, image["path"], image["title"], layout_index=1)  # Adjust layout_index if needed

# Save the modified presentation
prs.save("output_presentation.pptx")
print("PowerPoint presentation created successfully using the template!")